#!/usr/bin/env python3.12
"""Build an Aecon-branded PPTX deck from scratch using python-pptx.
Reference pattern: aecon-enclave-procurement-deck-2026-07-20.pptx

Usage:
  python3.12 build_aecon_pptx.py

This is a TEMPLATE — copy and modify for specific decks.
Key principles:
  - 13.333" x 7.5" widescreen
  - Arial font (Univers not available on recipient machines)
  - Print red #E51937 for PPTX
  - MSO_SHAPE.ROUNDED_RECTANGLE cards
  - Thin red bar takeaway boxes
  - Colored RACI cells
  - Blank layout (no built-in placeholders)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══ BRAND ═══
RED = RGBColor(0xE5, 0x19, 0x37)
CHARCOAL = RGBColor(0x25, 0x25, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GRAY = RGBColor(0x46, 0x46, 0x46)
SILVER = RGBColor(0x74, 0x76, 0x79)
MID_GRAY = RGBColor(0x70, 0x70, 0x70)
BORDER = RGBColor(0xEA, 0xEA, 0xEA)
C0 = RGBColor(0xC0, 0xC0, 0xC0)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DARK_RED = RGBColor(0xEF, 0x44, 0x44)

LOGO = '/data/nextcloud/data/amyn/files/briefings/aecon-assets/logo-aecon-red.png'
OUT = '/data/nextcloud/data/amyn/files/briefings/DECK_NAME.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ═══ HELPERS ═══

def tb(slide, l, t, w, h, text="", fs=12, bold=False, color=BODY_GRAY, align=PP_ALIGN.LEFT, font='Arial'):
    """Add a single-run textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    p.alignment = align
    return tf

def rich_tb(slide, l, t, w, h):
    """Add a textbox for manual paragraph/run building."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    return tf

def arun(tf, text, fs=12, bold=False, color=BODY_GRAY, font='Arial', align=PP_ALIGN.LEFT):
    """Add a paragraph+run to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return r

def bg(slide, color):
    """Set slide background color."""
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def footer(slide, left_text, right_text, dark=False):
    """Add slide footer."""
    c = MID_GRAY if dark else SILVER
    tb(slide, 0.5, 7.05, 6, 0.3, left_text, 7, False, c)
    tb(slide, 7.3, 7.05, 5.5, 0.3, right_text, 7, False, c, PP_ALIGN.RIGHT)

def snum(slide, n, total=14, dark=False):
    """Add slide number badge."""
    c = MID_GRAY if dark else SILVER
    tb(slide, 12.3, 0.12, 0.8, 0.25, f"{n} / {total}", 7, True, c, PP_ALIGN.RIGHT)

def logo(slide, l=0.5, t=0.15, w=1.1):
    """Embed Aecon logo."""
    slide.shapes.add_picture(LOGO, Inches(l), Inches(t), Inches(w))

def card(slide, l, t, w, h, title="", body="", border_color=None):
    """Add a rounded-rect card with title and body text."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = border_color if border_color else BORDER
    s.line.width = Pt(1.5) if border_color else Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    if title:
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = CHARCOAL; r.font.name = 'Arial'
    if body:
        p2 = tf.add_paragraph() if title else tf.paragraphs[0]
        p2.space_before = Pt(4)
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(10); r2.font.color.rgb = BODY_GRAY; r2.font.name = 'Arial'
    return s

def takeaway(slide, text, top=6.2):
    """Add red-left-border takeaway box."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top), Inches(0.03), Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    tf = tb(slide, 0.7, top - 0.01, 12, 0.5, "", 10)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Bottom line: "; r1.font.size = Pt(10); r1.font.bold = True
    r1.font.color.rgb = RED; r1.font.name = 'Arial'
    r2 = p.add_run(); r2.text = text; r2.font.size = Pt(10)
    r2.font.color.rgb = CHARCOAL; r2.font.name = 'Arial'

def impact(slide, l, t, w, num, label, desc):
    """Big number stat box."""
    tb(slide, l, t, w, 0.7, num, 38, True, RED, PP_ALIGN.CENTER)
    tb(slide, l, t+0.65, w, 0.25, label, 12, True, CHARCOAL, PP_ALIGN.CENTER)
    tb(slide, l, t+0.95, w, 0.5, desc, 9, False, BODY_GRAY, PP_ALIGN.CENTER)

def action_title(slide, parts, top=0.25):
    """Rich action title. parts = [(text, is_red), ...]"""
    tf = rich_tb(slide, 0.5, top, 12.3, 1.0)
    p = tf.paragraphs[0]
    for text, is_red in parts:
        r = p.add_run(); r.text = text
        r.font.size = Pt(24); r.font.bold = True; r.font.name = 'Arial'
        r.font.color.rgb = RED if is_red else CHARCOAL

def step_card(slide, l, t, w, num, name, owner, tint=False):
    """Process flow step card."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(1.1))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xFE, 0xF5, 0xF5) if tint else WHITE
    s.line.color.rgb = BORDER; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08); tf.margin_top = Inches(0.06)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = RED; r.font.name = 'Arial'
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = name; r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = CHARCOAL; r2.font.name = 'Arial'
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = owner; r3.font.size = Pt(8); r3.font.color.rgb = SILVER; r3.font.name = 'Arial'

def tl_item(slide, top, phase, title, desc, hours):
    """Timeline item with dot."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(top + 0.02), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = CHARCOAL; dot.line.fill.background()
    tb(slide, 0.72, top - 0.03, 1.5, 0.2, phase, 9, True, RED)
    tb(slide, 0.72, top + 0.18, 4, 0.2, title, 12, True, CHARCOAL)
    tb(slide, 0.72, top + 0.4, 5.5, 0.4, desc, 9, False, BODY_GRAY)
    tb(slide, 0.72, top + 0.75, 1, 0.2, hours, 10, True, RED)

def raci_cell(slide, l, t, letter, color):
    """Small colored RACI badge."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(0.35), Inches(0.28))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = letter
    r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = WHITE if color not in [RGBColor(0xE8,0xE8,0xE8), RGBColor(0xF5,0xF5,0xF5)] else BODY_GRAY
    r.font.name = 'Arial'

# ═══ BUILD DECK ═══
# Add slides below using the helpers above.
# See fcs-access-clearance-leadership-deck.pptx for a complete 14-slide example.

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Size: {os.path.getsize(OUT):,} bytes")
print(f"Slides: {len(prs.slides)}")
